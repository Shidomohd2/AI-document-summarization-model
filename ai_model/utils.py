# utils.py
import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document as DocxDocument
import pytesseract
from PIL import Image
import logging
from pdf2image import convert_from_path
import re
from transformers import MBartForConditionalGeneration, MBart50TokenizerFast

logger = logging.getLogger(__name__)

# Load the model and tokenizer once

model_name = "facebook/mbart-large-50"
tokenizer = MBart50TokenizerFast.from_pretrained(model_name)
model = MBartForConditionalGeneration.from_pretrained(model_name)

def clean_text(text):
    """Clean and preprocess input text."""
    # Remove extra spaces
    text = re.sub(r'\s+', ' ', text)
    
    # Remove special characters (keep only letters, numbers, and basic punctuation)
    text = re.sub(r'[^\w\s.,!?]', '', text)
    
    # Remove repeated phrases
    words = text.split()
    unique_words = []
    for word in words:
        if word not in unique_words[-5:]:  # Avoid adding the same word repeatedly
            unique_words.append(word)
    text = " ".join(unique_words)
    
    return text.strip()

def extract_text_from_image(image):
    """Extract text from an image using OCR."""
    try:
        return pytesseract.image_to_string(image, lang='ara')
    except Exception as e:
        logger.error(f"OCR Error: {e}")
        return ""

def extract_images_from_pdf(file_path):
    """Extract text from images embedded in a PDF."""
    text = ""
    try:
        images = convert_from_path(file_path)
        for image in images:
            text += extract_text_from_image(image) + "\n"
    except Exception as e:
        logger.error(f"PDF Image Error: {e}")
    return text

def extract_text_from_pdf(file_path):
    """Extract text from a PDF file."""
    text = ""
    try:
        with fitz.open(file_path) as doc:
            for page in doc:
                text += page.get_text() + "\n\n"  # Double newline between pages
        text += extract_images_from_pdf(file_path)
    except Exception as e:
        logger.error(f"PDF Error: {e}")
    return text.strip()

def extract_text_from_txt(file_path):
    """Extract text from a TXT file."""
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()
    except Exception as e:
        logger.error(f"TXT Error: {e}")
        return ""

def extract_text_from_pptx(file_path):
    """Extract text from a PPTX file."""
    text = ""
    try:
        presentation = Presentation(file_path)
        for slide in presentation.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            text += "\n".join(slide_text) + "\n\n"
    except Exception as e:
        logger.error(f"PPTX Error: {e}")
    return text.strip()

def extract_text_from_docx(file_path):
    """Extract text from a DOCX file."""
    text = []
    try:
        doc = DocxDocument(file_path)
        for paragraph in doc.paragraphs:
            text.append(paragraph.text)
        return "\n\n".join(text)
    except Exception as e:
        logger.error(f"DOCX Error: {e}")
        return ""

def summarize_text(text):
    """Summarize text using the loaded model."""
    try:
        text = clean_text(text)  # Clean the input text
        logger.info(f"Input text for summarization: {text[:500]}...")  # Log first 500 characters
        
        inputs = TOKENIZER(text, return_tensors="pt", max_length=1024, truncation=True)
        summary_ids = MODEL.generate(
            inputs["input_ids"],
            max_length=150,
            min_length=30,
            length_penalty=2.0,
            num_beams=4,
            repetition_penalty=2.5,
            no_repeat_ngram_size=3,
        )
        summary = TOKENIZER.decode(summary_ids[0], skip_special_tokens=True)
        logger.info(f"Generated summary: {summary}")  # Log the summary
        return summary
    except Exception as e:
        logger.error(f"Summarization Error: {e}")
        return "Error generating summary."

def extract_pages(file_path):
    """Extract pages from a file and return them as a list of strings."""
    pages = []
    try:
        if file_path.endswith('.pdf'):
            with fitz.open(file_path) as doc:
                for page in doc:
                    pages.append(page.get_text())
        elif file_path.endswith('.pptx'):
            presentation = Presentation(file_path)
            for slide in presentation.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                pages.append("\n".join(slide_text))
        elif file_path.endswith('.docx'):
            doc = DocxDocument(file_path)
            for paragraph in doc.paragraphs:
                pages.append(paragraph.text)
        elif file_path.endswith('.txt'):
            with open(file_path, 'r', encoding='utf-8') as file:
                pages = file.read().split('\n\n')  # Split by paragraphs
        else:
            raise ValueError("Unsupported file format")
    except Exception as e:
        logger.error(f"Page extraction error: {e}")
    return pages

def process_file(file_content):
    """Process file content and return a summary or cleaned text."""
    try:
        cleaned_text = clean_text(file_content)
        
        # Skip summarization if the text is too short or repetitive
        if len(cleaned_text.split()) < 20:  # Adjust threshold as needed
            return cleaned_text  # Return the cleaned text as-is
        
        return summarize_text(cleaned_text)
    except Exception as e:
        logger.error(f"Processing Error: {e}")
        return f"Error: {str(e)}"